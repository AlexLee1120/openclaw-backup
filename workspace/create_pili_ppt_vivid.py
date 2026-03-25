import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_pili_ppt_with_images(output_path):
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "霹靂國際多媒體戰略分析報告"
    subtitle.text = "傳統布袋戲的「大轉骨」與泛娛樂 IP 帝國的戰略重塑\n報告人：小艾 (Digital CTO)"

    def add_bullet_slide(title_text, bullets, image_path=None):
        if image_path:
            # Layout 8 is 'Picture with Caption' but let's just use layout 1 and adjust manually
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = title_text
            
            # Left side for text
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
            
            # Shrink text box width to 50%
            body_shape.width = Inches(4.5)
            
            # Add image on right
            try:
                slide.shapes.add_picture(image_path, Inches(5), Inches(1.5), height=Inches(5))
            except Exception as e:
                print(f"Error adding image {image_path}: {e}")
        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = title_text
            tf = slide.placeholders[1].text_frame
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    # Slide 2: Q-a (1/2)
    add_bullet_slide("a. 核心資源與競爭能力 (1/2)", [
        "百年家族傳承與品牌威信：雲林黃家四代技藝傳承，確立「正宗」地位與戲迷忠誠度。",
        "「一條龍」垂直整合生產模式：從編劇、口白到後製、特效，保證風格高度一致性。",
        "極高競爭壁壘：全流程內部完成，累積深厚技術經驗值，難以被對手複製。"
    ])

    # Slide 3: Q-a (2/2) - Tech
    add_bullet_slide("a. 核心資源與競爭能力 (2/2)", [
        "技術革命：領先導入 4K/8K 攝影、3D 列印戲偶、BJD 技術與 AI 語音模擬。",
        "數位資產化與 IP 確權：建立數位典藏系統，將 4,000 多名角色符號化與圖騰化。",
        "IP 生命週期極大化：透過數位確權，快速轉化為手遊、周邊與 NFT 資產。"
    ], "/Users/alex-ai/.openclaw/workspace/ppt_tech_workshop.png")

    # Slide 4: Q-b (1/2) - IP Ecosystem
    add_bullet_slide("b. 經營模式分析", [
        "核心戰略：一源多用 (One Source, Multi-use) 的 IP 生態系。",
        "影視發行轉型：DVD 出租轉向自有的「PILI 線上看」串流平台 (OTT)。",
        "宇宙重啟 (Reboot)：透過《素還真》大電影重啟英雄宇宙，解決長壽劇敘事負擔。",
        "多角化經營：手遊、桌遊、Web3 (NFT) 及跨國品牌合作 (如暴雪娛樂)。"
    ], "/Users/alex-ai/.openclaw/workspace/ppt_ip_ecosystem.png")

    # Slide 5: Q-b (2/2)
    add_bullet_slide("b. 戰略評論與挑戰", [
        "模式優勢：資產輕量化潛力，數位化降低對個別操偶大師的依賴。",
        "轉型陣痛：垂直整合模式帶來高昂研發與人力成本，導致短期財務壓力 (如 2023 年)。",
        "關鍵轉向：從「封閉系統」轉向「開放式協作 (製作委員會模式)」，以輕資產全球運作。"
    ])

    # Slide 6: Q-c (1/2)
    add_bullet_slide("c. 海外擴張策略：製作委員會", [
        "商業模式突破：放棄單純授權，改採日式「製作委員會」制度 (利益共享、風險共擔)。",
        "深度捆綁合作：與 Nitro+、Good Smile Company 等日本指標企業建立戰略同盟。",
        "標竿案例：2016 年啟動之《東離劍遊紀》，成為台灣文化輸出國際的成功範式。"
    ])

    # Slide 7: Q-c (2/2)
    add_bullet_slide("c. 文化翻譯與門檻降低", [
        "敘事重寫：邀請虛淵玄重新編寫王道冒險劇本，簡化漢學禪道門檻。",
        "美學折衷：戲偶造型針對日系審美進行「動漫化/萌化」，提升視覺吸引力。",
        "策略核心：在敘事權上跨國合作，保留核心操偶技術作為靈魂，「以退為進」。"
    ])

    # Slide 8: Q-d (1/2)
    add_bullet_slide("d. 三年戰略規劃 (2028)", [
        "財務目標：營考挑戰 12 億台幣，EPS 重返 2.5 元以上。",
        "技術降本：全面落實 AI 語音生成與自動化後製，降低單集成本 20% 以上。",
        "核心產品：穩定產出「英雄重啟」系列大片與跨國合製劇集。",
        "社群經營：修補新舊戲迷衝突感，深耕 OTT 高黏著度訂閱基礎。"
    ])

    # Slide 9: Q-d (2/2) - Metaverse
    add_bullet_slide("d. 十年願景規劃 (2035)", [
        "財務目標：營收挑戰 50-80 億台幣，海外營收佔比突破 70%。",
        "技術願景：建立「霹靂 XR 元宇宙」，實現全球沉浸式觀影與虛擬互動。",
        "內容發布：AI 驅動敘事引擎，實現內容在全球市場的「零時差」同步發布。",
        "永續價值：落實 ESG 淨零排放，成為全球傳統文化科技化的永續標竿。"
    ], "/Users/alex-ai/.openclaw/workspace/ppt_xr_theater.png")

    # Slide 10: Ending Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "結語：Keep on Rolling"
    content = slide.placeholders[1].text_frame
    content.text = "傳統技藝的重生，在於不間斷的自我轉化與技術賦能。"

    prs.save(output_path)
    print(f"PPT with images created successfully at {output_path}")

if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else "/Users/alex-ai/.openclaw/workspace/霹靂國際多媒體戰略分析簡報_精裝插圖版.pptx"
    create_pili_ppt_with_images(path)
